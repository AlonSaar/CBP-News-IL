"""
Generates a PowerPoint presentation (.pptx) from the articles store.
Output: output/cbp_hebrew_reviews.pptx
"""

from __future__ import annotations

import io
import logging
from datetime import datetime
from pathlib import Path
from typing import Optional

import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

import config

logger = logging.getLogger(__name__)

OUTPUT_DIR = config.ROOT / "output"
PPTX_PATH = OUTPUT_DIR / "cbp_hebrew_reviews.pptx"

C_NAVY  = RGBColor(0x1a, 0x1a, 0x2e)
C_RED   = RGBColor(0xe9, 0x45, 0x60)
C_BLUE  = RGBColor(0x4c, 0x5f, 0xd5)
C_GRAY  = RGBColor(0x88, 0x88, 0x88)
C_DARK  = RGBColor(0x22, 0x22, 0x33)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT = RGBColor(0xF0, 0xF0, 0xF5)

SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)


# ---------------------------------------------------------------------------
# Core RTL text helper
# ---------------------------------------------------------------------------

def _set_para_rtl(para):
    """Set paragraph to RTL + right-align in the XML."""
    pPr = para._p.get_or_add_pPr()
    pPr.set("algn", "r")
    # Remove existing rtl elements then re-add
    for el in pPr.findall(qn("a:rtl")):
        pPr.remove(el)
    rtl_el = etree.SubElement(pPr, qn("a:rtl"))
    rtl_el.set("val", "1")


def _set_run_lang(run, lang="he-IL"):
    """Set language hint on a run so PowerPoint handles bidi correctly."""
    rPr = run._r.get_or_add_rPr()
    rPr.set("lang", lang)
    rPr.set("altLang", "en-US")


def _add_solid_rect(slide, x, y, w, h, color: RGBColor):
    """Add a filled rectangle with no border."""
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_text(slide, x, y, w, h, text, font_size, bold=False,
              color=C_DARK, rtl=True, italic=False, align=PP_ALIGN.RIGHT):
    """Add a single-paragraph text box."""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    if rtl:
        _set_para_rtl(p)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Arial"
    if rtl:
        _set_run_lang(run)
    return tf


def _add_body_text(slide, x, y, w, h, body, font_size, color=C_DARK):
    """
    Add Hebrew body text as multiple paragraphs — split on '. '
    so each sentence gets its own line with proper RTL alignment.
    """
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Split into sentences on period + space
    sentences = []
    for chunk in body.split(". "):
        chunk = chunk.strip()
        if chunk:
            sentences.append(chunk)

    for i, sentence in enumerate(sentences):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.RIGHT
        _set_para_rtl(p)

        # Add period back (except possibly last if body ended without one)
        display = sentence if sentence.endswith(".") else sentence + "."
        run = p.add_run()
        run.text = display
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Arial"
        _set_run_lang(run)

        # Add spacing after each sentence
        pPr = p._p.get_or_add_pPr()
        spcAft = etree.SubElement(pPr, qn("a:spcAft"))
        spcPts = etree.SubElement(spcAft, qn("a:spcPts"))
        spcPts.set("val", "600")  # 6pt spacing after each sentence

    return tf


def _download_image(url: str) -> Optional[bytes]:
    """Download image and convert WEBP → JPEG so python-pptx can embed it."""
    try:
        resp = requests.get(url, timeout=10,
                            headers={"User-Agent": config.CBP_USER_AGENT})
        resp.raise_for_status()
        data = resp.content

        # Convert WEBP (and any other unsupported format) to JPEG via Pillow
        content_type = resp.headers.get("Content-Type", "").lower()
        if "webp" in content_type or url.lower().endswith(".webp") or b"WEBP" in data[:20]:
            try:
                from PIL import Image as PILImage
                buf_in = io.BytesIO(data)
                buf_out = io.BytesIO()
                img = PILImage.open(buf_in).convert("RGB")
                img.save(buf_out, format="JPEG", quality=90)
                data = buf_out.getvalue()
                logger.debug("Converted WEBP → JPEG (%d bytes)", len(data))
            except Exception as conv_exc:
                logger.warning("WEBP conversion failed: %s", conv_exc)
                return None

        return data
    except Exception as exc:
        logger.warning("Could not download image %s: %s", url, exc)
        return None


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def _add_cover_slide(prs: Presentation, articles: list[dict]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Full dark background
    _add_solid_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, C_NAVY)

    # Top accent strip
    _add_solid_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.12), C_RED)

    # Bottom accent strip
    _add_solid_rect(slide, Inches(0), Inches(5.5), SLIDE_W, Inches(0.125), C_RED)

    # Main Hebrew title
    _add_text(slide, Inches(0.5), Inches(1.4), Inches(9), Inches(1.0),
              "סקירות CBP בעברית", 40, bold=True, color=C_WHITE, rtl=True)

    # Red divider line
    _add_solid_rect(slide, Inches(0.5), Inches(2.55), Inches(9), Inches(0.05), C_RED)

    # Subtitle
    _add_text(slide, Inches(0.5), Inches(2.7), Inches(9), Inches(0.6),
              "תפיסות סמים, נשק, חיות ומוצרים אסורים", 18, color=C_RED, rtl=True)

    # Stats
    quarters = sorted(set(a.get("quarter", "") for a in articles), reverse=True)
    now_str = datetime.now().strftime("%d/%m/%Y")
    _add_text(slide, Inches(0.5), Inches(3.4), Inches(9), Inches(0.5),
              f"עודכן: {now_str}   |   כתבות: {len(articles)}   |   {', '.join(quarters)}",
              12, color=C_GRAY, rtl=True)

    # Footer
    _add_text(slide, Inches(0.5), Inches(5.1), Inches(9), Inches(0.35),
              "נציגות ניו יורק  |  רשות המסים", 10, color=C_GRAY, rtl=True)


def _add_quarter_divider(prs: Presentation, quarter: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_solid_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, RGBColor(0x10, 0x10, 0x22))
    _add_solid_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.1), C_RED)
    _add_solid_rect(slide, Inches(0), Inches(5.52), SLIDE_W, Inches(0.1), C_RED)

    _add_text(slide, Inches(1), Inches(1.9), Inches(8), Inches(1.0),
              quarter, 52, bold=True, color=C_WHITE, rtl=False, align=PP_ALIGN.CENTER)

    names = {"Q1": "ינואר – מרץ", "Q2": "אפריל – יוני",
             "Q3": "יולי – ספטמבר", "Q4": "אוקטובר – דצמבר"}
    q_num = quarter.split("-")[0] if "-" in quarter else quarter
    year = quarter.split("-")[1] if "-" in quarter else ""
    label = f"{names.get(q_num, '')} {year}".strip()
    if label:
        _add_text(slide, Inches(1), Inches(3.1), Inches(8), Inches(0.6),
                  label, 22, color=C_RED, rtl=True, align=PP_ALIGN.CENTER)


def _add_article_slide(prs: Presentation, art: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # White background
    _add_solid_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, C_WHITE)

    # Header bar
    _add_solid_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.52), C_NAVY)

    # Date + quarter in header (LTR, left side)
    date_q = f"{art.get('article_date', '')}   |   {art.get('quarter', '')}"
    _add_text(slide, Inches(0.2), Inches(0.08), Inches(9.6), Inches(0.38),
              date_q, 9, color=C_GRAY, rtl=False, align=PP_ALIGN.LEFT)

    # Image (right column)
    IMG_X, IMG_Y, IMG_W, IMG_H = Inches(6.4), Inches(0.62), Inches(3.4), Inches(3.5)
    image_url = art.get("image_url", "")
    img_added = False
    if image_url:
        img_bytes = _download_image(image_url)
        if img_bytes:
            try:
                slide.shapes.add_picture(io.BytesIO(img_bytes), IMG_X, IMG_Y, IMG_W, IMG_H)
                img_added = True
            except Exception as exc:
                logger.warning("Image insert failed: %s", exc)
    if not img_added:
        _add_solid_rect(slide, IMG_X, IMG_Y, IMG_W, IMG_H, C_LIGHT)
        _add_text(slide, IMG_X, Inches(2.1), IMG_W, Inches(0.5),
                  "CBP", 18, color=C_GRAY, rtl=False, align=PP_ALIGN.CENTER)

    # Left text column
    TX, TW = Inches(0.25), Inches(5.9)

    # Title (large, bold)
    _add_text(slide, TX, Inches(0.62), TW, Inches(1.05),
              art.get("title", ""), 15, bold=True, color=C_NAVY, rtl=True)

    # Thin red line under title
    _add_solid_rect(slide, TX, Inches(1.72), TW, Inches(0.025), C_RED)

    # Location
    _add_text(slide, TX, Inches(1.78), TW, Inches(0.38),
              art.get("location", ""), 10, bold=True, color=C_RED, rtl=True)

    # Crossing type
    ct = art.get("crossing_type", "").replace("סוג מעבר:", "").strip()
    if ct:
        _add_text(slide, TX, Inches(2.18), TW, Inches(0.35),
                  "סוג מעבר: " + ct, 9, italic=True, color=C_BLUE, rtl=True)

    # Body — split into sentences with paragraph spacing
    body = art.get("body", "")
    if body:
        _add_body_text(slide, TX, Inches(2.58), TW, Inches(2.8),
                       body, 9.5, color=C_DARK)

    # Footer bar
    _add_solid_rect(slide, Inches(0), Inches(5.28), SLIDE_W, Inches(0.345), C_LIGHT)
    url = art.get("url", "")
    if url:
        _add_text(slide, Inches(0.2), Inches(5.3), Inches(9.6), Inches(0.28),
                  f"מקור: {url}", 7, color=C_GRAY, rtl=False, align=PP_ALIGN.LEFT)


# ---------------------------------------------------------------------------
# Main entry point
# ---------------------------------------------------------------------------

def generate(articles: list[dict]) -> Path:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _add_cover_slide(prs, articles)

    quarters: dict[str, list[dict]] = {}
    for art in articles:
        q = art.get("quarter", "Q?-????")
        quarters.setdefault(q, []).append(art)

    for q in sorted(quarters.keys(), reverse=True):
        _add_quarter_divider(prs, q)
        for art in sorted(quarters[q], key=lambda a: a.get("article_date", ""), reverse=True):
            _add_article_slide(prs, art)

    prs.save(PPTX_PATH)
    logger.info("PowerPoint written: %s (%d articles)", PPTX_PATH, len(articles))
    return PPTX_PATH
